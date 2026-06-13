from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.domain.models import FileDescriptor, ParsedDocument, ParsedSection
from app.services.parsers.assets import write_image_asset
from app.services.parsers.base import BaseParser


class PPTXParser(BaseParser):
    def parse(self, file_descriptor: FileDescriptor) -> ParsedDocument:
        from pptx import Presentation

        presentation = Presentation(str(file_descriptor.file_path))
        smart_art_shape_type = getattr(MSO_SHAPE_TYPE, "SMART_ART", None)
        diagram_shape_type = getattr(MSO_SHAPE_TYPE, "DIAGRAM", None)

        sections: list[ParsedSection] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_text_parts: list[str] = []
            table_parts: list[str] = []
            visual_parts: list[str] = []
            visual_refs: list[dict[str, str]] = []

            for shape_index, shape in enumerate(slide.shapes, start=1):
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_text_parts.append(text.strip())

                if getattr(shape, "has_table", False):
                    table = shape.table
                    for row in table.rows:
                        row_values = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        if any(row_values):
                            table_parts.append(" | ".join(row_values))

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    visual_parts.append(f"[IMAGE] name={shape.name}")
                    image = getattr(shape, "image", None)
                    if image:
                        image_path = write_image_asset(
                            file_path=file_descriptor.file_path,
                            blob=image.blob,
                            ext=image.ext,
                            name=f"slide_{index}_shape_{shape_index}",
                        )
                        visual_refs.append(
                            {
                                "type": "image",
                                "name": shape.name,
                                "path": image_path,
                                "slide_number": str(index),
                            }
                        )
                elif shape.shape_type in {
                    smart_art_shape_type,
                    diagram_shape_type,
                    MSO_SHAPE_TYPE.GROUP,
                    MSO_SHAPE_TYPE.FREEFORM,
                    MSO_SHAPE_TYPE.AUTO_SHAPE,
                }:
                    # These shape classes frequently represent architecture diagrams.
                    if not text or not text.strip():
                        visual_parts.append(f"[DIAGRAM_SHAPE] type={shape.shape_type} name={shape.name}")
                        visual_refs.append(
                            {
                                "type": "diagram",
                                "name": shape.name,
                                "shape_type": str(shape.shape_type),
                                "slide_number": str(index),
                            }
                        )

            slide_text = "\n".join(slide_text_parts).strip()
            table_text = "\n".join(table_parts).strip()
            visual_text = "\n".join(dict.fromkeys(visual_parts)).strip()

            if slide_text:
                sections.append(
                    ParsedSection(
                        text=slide_text,
                        slide_number=index,
                        metadata={"visual_refs": visual_refs},
                    )
                )
            if table_text:
                sections.append(
                    ParsedSection(
                        text="[TABLE]\n" + table_text,
                        slide_number=index,
                        section_title=f"Slide {index} tables",
                        metadata={"visual_refs": visual_refs},
                    )
                )
            if visual_text:
                sections.append(
                    ParsedSection(
                        text=visual_text,
                        slide_number=index,
                        section_title=f"Slide {index} visuals",
                        metadata={"visual_refs": visual_refs},
                    )
                )

        return ParsedDocument(document_id=file_descriptor.document_id, sections=sections)
